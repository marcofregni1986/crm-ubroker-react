import collections
import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import json

def extract_content(pptx_path, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    img_dir = os.path.join(output_dir, 'slide_images')
    if not os.path.exists(img_dir):
        os.makedirs(img_dir)

    try:
        prs = Presentation(pptx_path)
        slides_data = []

        for i, slide in enumerate(prs.slides):
            slide_info = {
                "number": i + 1,
                "title": "",
                "text_blocks": [],
                "images": [],
                "tables": []
            }

            # Try to get title
            if slide.shapes.title:
                slide_info["title"] = slide.shapes.title.text.strip()

            def process_shapes(shapes):
                for shape in shapes:
                    # Recursive for groups
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        process_shapes(shape.shapes)
                    
                    # Text extraction
                    if hasattr(shape, "text") and shape.text.strip():
                        # Skip if it's the title (already got it)
                        if slide.shapes.title and shape == slide.shapes.title:
                            continue
                        txt = shape.text.strip()
                        # Clean encoding artifacts
                        txt = txt.replace('Þ', 'è').replace('Ó', 'à').replace('ò', '·').replace('Æ', '’').replace('░', 'ª').replace('└', 'À').replace('àe', ' e').replace('à', '...')
                        slide_info["text_blocks"].append(txt)
                    
                    # Image extraction
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Try standard way
                            image = shape.image
                            image_filename = f"slide_{i+1}_img_{len(slide_info['images'])+1}.{image.ext}"
                            image_path = os.path.join(img_dir, image_filename)
                            with open(image_path, 'wb') as f:
                                f.write(image.blob)
                            slide_info["images"].append(f"slide_images/{image_filename}")
                        except Exception as e_img:
                            print(f"Skipping an image on slide {i+1} due to error: {e_img}")
                    
                    # Table extraction
                    if shape.has_table:
                        table_rows = []
                        for row in shape.table.rows:
                            cells = [cell.text_frame.text.strip() for cell in row.cells]
                            table_rows.append(" | ".join(cells))
                        slide_info["tables"].append(table_rows)

            process_shapes(slide.shapes)
            
            # Determine type
            if not slide_info["title"] and len(slide_info["text_blocks"]) <= 2 and not slide_info["images"]:
                slide_info["type"] = "title"
                if slide_info["text_blocks"]:
                    slide_info["title"] = slide_info["text_blocks"][0]
                    slide_info["subtitle"] = slide_info["text_blocks"][1] if len(slide_info["text_blocks"]) > 1 else ""
                    slide_info["text_blocks"] = []
            else:
                slide_info["type"] = "content"

            slides_data.append(slide_info)

        # Save JSON
        json_path = os.path.join(output_dir, 'slides_full.json')
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(slides_data, f, indent=4, ensure_ascii=False)
            
        print(f"Extracted {len(slides_data)} slides. Images saved to {img_dir}")
        return slides_data

    except Exception as e:
        print(f"Error during extraction: {e}")
        return None

if __name__ == "__main__":
    pptx_path = r"C:\Stepone\pp NUOVO STEP ONE  2024 2.pptx"
    output_dir = r"c:\modifica crm rise"
    extract_content(pptx_path, output_dir)
