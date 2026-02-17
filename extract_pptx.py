import collections
import collections.abc
from pptx import Presentation
import sys
import os

def extract_text_from_pptx(file_path):
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        return

    try:
        prs = Presentation(file_path)
        print(f"--- Extraction from: {os.path.basename(file_path)} ---")
        print(f"Total slides: {len(prs.slides)}")
        
        for i, slide in enumerate(prs.slides):
            print(f"\nSlide {i + 1}:")
            
            # Try to get title shape specifically if it exists
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
                if title:
                    print(f"TITLE: {title}")
            
            for shape in slide.shapes:
                # Skip the title shape as we already printed it
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                    
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        print(f"TEXT: {text}")
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text_frame.text.strip() for cell in row.cells]
                        print(f"TABLE: {' | '.join(row_text)}")
        
        print("\n--- End of Extraction ---")
    except Exception as e:
        print(f"An error occurred: {e}")

if __name__ == "__main__":
    pptx_path = r"C:\Stepone\pp NUOVO STEP ONE  2024 2.pptx"
    extract_text_from_pptx(pptx_path)
